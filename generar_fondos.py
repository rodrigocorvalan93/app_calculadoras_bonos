# -*- coding: utf-8 -*-
# generar_fondos.py
#
# Lee RF_Detalle_Carteras.xlsm y genera:
#   1. Slides de fondos embebibles en el HTML del comité
#   2. Una presentación PPTX separada
#
# USO:
#   from generar_fondos import generar_fondos_html, generar_fondos_pptx
#
#   html_slides = generar_fondos_html(
#       rf_detalle_path = r"C:\...\RF_Detalle_Carteras.xlsm"
#   )
#
#   generar_fondos_pptx(
#       rf_detalle_path = r"C:\...\RF_Detalle_Carteras.xlsm",
#       output_path     = r"C:\...\comite_fondos.pptx",
#       template_path   = r"C:\...\Presentación_Comité_Claude.pptx",  # opcional
#   )
# =============================================================================

from __future__ import annotations

import json
import os
import re
from datetime import date
from pathlib import Path
from typing import Optional

import numpy as np
import pandas as pd

# =============================================================================
# CONFIGURACIÓN DE FONDOS
# =============================================================================

# Cada fondo define:
#   sheet_rf:  sheet en RF_Detalle_Carteras.xlsm
#   col_metr:  columna (0-indexed) donde están los valores de métricas (TIR, Duration, Patrimonio)
#   col_pct:   columna (0-indexed) del % de participación de cada instrumento
#   col_monto: columna (0-indexed) del monto nominal
#   nombre:    nombre para mostrar
#   tipo_vcp:  título del gráfico de variaciones VCP
#   tiene_tir: True si muestra TIR en KPIs (Multimercado no la muestra en la slide)

FONDOS = [
    {
        "id":         "performance",
        "nombre":     "Delta Performance",
        "sheet_rf":   "Cash Management",
        "col_metr":   13,   # PERFORMANCE
        "col_pct":    14,
        "col_monto":  None,
        "tiene_tir":  True,
        "tipo_vcp":   "Renta Fija T+0",
        "peers_label":"Performance",
    },
    {
        "id":         "ahorro_plus",
        "nombre":     "Delta Ahorro Plus",
        "sheet_rf":   "Cash Management",
        "col_metr":   10,   # AHORRO PLUS
        "col_pct":    11,
        "col_monto":  None,
        "tiene_tir":  True,
        "tipo_vcp":   'T+1 "Agresivo"',
        "peers_label":"Ahorro Plus",
    },
    {
        "id":         "retorno_real",
        "nombre":     "Delta Retorno Real",
        "sheet_rf":   "CER",
        "col_metr":   7,
        "col_pct":    8,
        "col_monto":  None,
        "tiene_tir":  True,
        "tipo_vcp":   "CER (Inflación)",
        "peers_label":"Delta Retorno Real",
    },
    {
        "id":         "multimercado",
        "nombre":     "Delta Multimercado I",
        "sheet_rf":   "Largo Plazo",
        "col_metr":   16,
        "col_pct":    17,
        "col_monto":  None,
        "tiene_tir":  False,
        "tipo_vcp":   "Total Return (Renta Mixta)",
        "peers_label":"Multimercado I",
    },
]

# Mapeo subcategoría RF_Detalle → tipo de ajuste para gráfico
AJUSTE_MAP = {
    "ARS Fija":               "Tasa Fija",
    "ARS BADLAR":             "Tasa Fija",
    "ARS TAMAR":              "TAMAR",
    "CER Fija":               "CER/UVA",
    "UVA Fija":               "CER/UVA",
    "Dual (Fija/TAMAR) Fija": "Duales",
    "USD Fija":               "Hard-Dollar",
    "USD en Dólares":         "Hard-Dollar",
    "USB en Dólares":         "Hard-Dollar",
    "USD-Linked Fija":        "Dólar-Linked",
    "Acciones":               "Acciones",
    "RV":                     "Acciones",
    "Caja":                   "Caja",
    "Money Market":           "Caja",
}

# Mapeo (L1, L2) -> etiqueta en torta de Composición General
COMP_PIE_MAP = {
    # Soberanos
    ("Bonos Soberano",    "CER Fija"):               "Soberanos CER",
    ("Bonos Soberano",    "ARS Fija"):               "Soberanos Tasa Fija",
    ("Bonos Soberano",    "ARS TAMAR"):              "Soberanos TAMAR",
    ("Bonos Soberano",    "ARS BADLAR"):             "Soberanos Tasa Fija",
    ("Bonos Soberano",    "USD en Dolares"):         "Soberanos USD",
    ("Bonos Soberano",    "USD Fija"):               "Soberanos USD",
    ("Bonos Soberano",    "USB en Dolares"):         "Soberanos USD",
    ("Bonos Soberano",    "USD-Linked Fija"):        "Soberanos USD",
    ("Bonos Soberano",    "Dual (Fija/TAMAR) Fija"): "Soberanos Duales",
    ("Letras Soberano",   "CER Fija"):               "Soberanos CER",
    ("Letras Soberano",   "ARS Fija"):               "Soberanos Tasa Fija",
    ("Letras Soberano",   "ARS TAMAR"):              "Soberanos TAMAR",
    ("Letras Soberano",   "USD Fija"):               "Soberanos USD",
    ("Letras Soberano",   "UVA Fija"):               "Soberanos CER",
    # Subsoberanos
    ("Bonos Subsoberano", "CER Fija"):               "Subsoberanos CER",
    ("Bonos Subsoberano", "ARS Fija"):               "Subsoberanos Tasa Fija",
    ("Bonos Subsoberano", "ARS TAMAR"):              "Subsoberanos TAMAR",
    ("Bonos Subsoberano", "USD Fija"):               "Subsoberanos USD",
    ("Letras Subsoberano","ARS TAMAR"):              "Subsoberanos TAMAR",
    ("Letras Subsoberano","UVA Fija"):               "Subsoberanos CER",
    ("Letras Subsoberano","CER Fija"):               "Subsoberanos CER",
    # Corporativos
    ("Bonos Corporativo", "ARS Fija"):               "Plazos Fijos/Pagare",
    ("Bonos Corporativo", "ARS BADLAR"):             "Plazos Fijos/Pagare",
    ("Bonos Corporativo", "ARS CAUCION"):            "Plazos Fijos/Pagare",
    ("Bonos Corporativo", "ARS TAMAR"):              "FF TAMAR",
    ("Bonos Corporativo", "UVA Fija"):               "ON CER/UVA",
    ("Bonos Corporativo", "USD Fija"):               "ON USD",
    ("Bonos Corporativo", "USD-Linked Fija"):        "ON USD",
    ("Bonos Corporativo", "USB Fija"):               "ON USD",
    # Fideicomisos Financieros
    ("Fideicomisos Financieros", "en Pesos ARS TAMAR"):  "FF TAMAR",
    ("Fideicomisos Financieros", "en Pesos ARS BADLAR"): "FF TAMAR",
    ("Fideicomisos Financieros", "USD USD Fija"):         "ON USD",
    ("Fideicomisos Financieros", "USB NC"):               "ON USD",
    # FCI Money Markets -> Caja
    ("FCI Money Markets", "en Pesos"):               "Caja",
    # Acciones
    ("Local",             "Acciones"):               "Acciones",
}


# Colores para la torta de composicion general (paleta de la imagen de referencia)
COMP_PIE_COLORS = {
    "Soberanos CER":          "#4472C4",   # azul medio
    "Soberanos Tasa Fija":    "#1F3864",   # azul oscuro navy
    "Soberanos TAMAR":        "#70AD47",   # verde medio
    "Soberanos USD":          "#C55A11",   # naranja oscuro
    "Soberanos Duales":       "#7030A0",   # violeta
    "Subsoberanos CER":       "#5B9BD5",   # azul claro
    "Subsoberanos Tasa Fija": "#2E75B6",   # azul medio-oscuro
    "Subsoberanos TAMAR":     "#A9D18E",   # verde claro
    "Subsoberanos USD":       "#ED7D31",   # naranja
    "Plazos Fijos/Pagare":    "#7B5EA7",   # violeta/uva
    "FF TAMAR":               "#375623",   # verde oscuro
    "ON CER/UVA":             "#264478",   # azul profundo
    "ON USD":                 "#843C0C",   # marron
    "Acciones":               "#5A9E6F",   # verde esmeralda
    "Caja":                   "#203864",   # azul muy oscuro (casi negro)
}

# Categorías nivel 1 que van directo a Caja
CAJA_L1 = {"Caja y Equivalentes", "Caja", "Money Market", "Cuenta Corriente"}

# Colores por tipo de ajuste
AJUSTE_COLORS = {
    "Tasa Fija":    "#1e6fba",
    "CER/UVA":      "#4ea8e8",
    "TAMAR":        "#0f2557",
    "Duales":       "#6b93c4",
    "Hard-Dollar":  "#b85a1a",
    "Dólar-Linked": "#e8a020",
    "Acciones":     "#5a9e6f",
    "Caja":         "#8896a5",
}

AJUSTE_ORDER = ["Tasa Fija", "CER/UVA", "TAMAR", "Duales",
                "Hard-Dollar", "Dólar-Linked", "Acciones", "Caja"]

# =============================================================================
# LECTURA DE RF_DETALLE
# =============================================================================

def _read_sheet(path: str, sheet: str) -> pd.DataFrame:
    """Lee un sheet de Excel. Si el archivo está bloqueado (abierto en Excel),
    hace una copia temporal antes de leer."""
    import tempfile, shutil, os
    try:
        return pd.read_excel(path, sheet_name=sheet, header=None, engine="openpyxl")
    except PermissionError:
        # Archivo bloqueado por Excel — copiar a temp y leer desde ahí
        ext = os.path.splitext(path)[1]
        with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as tmp:
            tmp_path = tmp.name
        try:
            shutil.copy2(path, tmp_path)
            return pd.read_excel(tmp_path, sheet_name=sheet, header=None, engine="openpyxl")
        finally:
            try: os.unlink(tmp_path)
            except: pass


def _get_metrics(df: pd.DataFrame, col_metr: int) -> dict:
    """Extrae métricas clave del encabezado del sheet."""
    def _v(row):
        v = df.iloc[row, col_metr]
        return float(v) if not pd.isna(v) else None

    return {
        "duration":   _v(0),   # fila 0: Vida Prom
        "tir":        _v(1),   # fila 1: TIR bruta
        "fee":        _v(2),   # fila 2: Fee
        "tir_neta":   _v(3),   # fila 3: TIR neta
        "patrimonio": _v(5),   # fila 5: Patrimonio
    }


def _get_composition(df: pd.DataFrame, col_pct: int) -> dict:
    """
    Recorre el árbol de tenencias (nivel 1/2/4) y acumula % por tipo de ajuste.
    Retorna dict {tipo: pct_decimal}.
    """
    result = {}
    current_l2 = None

    for i in range(27, min(len(df), 300)):
        nivel = df.iloc[i, 0]
        cat   = str(df.iloc[i, 2]) if not pd.isna(df.iloc[i, 2]) else ""
        pct   = df.iloc[i, col_pct]

        if pd.isna(nivel):
            continue
        nivel = float(nivel)

        if nivel == 1.0:
            current_l2 = None
            # ¿Es una categoría de Caja directo?
            if cat in CAJA_L1 and not pd.isna(pct) and float(pct) > 1e-6:
                result["Caja"] = result.get("Caja", 0) + float(pct)

        elif nivel == 2.0:
            current_l2 = cat
            tipo = AJUSTE_MAP.get(cat)
            if tipo and not pd.isna(pct) and float(pct) > 1e-6:
                result[tipo] = result.get(tipo, 0) + float(pct)

        # nivel 4 — instrumento individual: no sumamos (ya está en nivel 2)

    # Limpiar valores muy pequeños
    result = {k: v for k, v in result.items() if v > 0.001}
    return result


def _get_composition_pie(df: pd.DataFrame, col_pct: int, col_monto: int,
                         patrimonio: float, liquidez: float) -> dict:
    """
    Recorre el arbol (L1, L2) y asigna cada combinacion a una etiqueta de torta
    usando COMP_PIE_MAP. Devuelve dict {etiqueta: pct} normalizado.
    Incluye Caja = liquidez / patrimonio.
    """
    result = {}

    # Caja desde liquidez
    if patrimonio and patrimonio > 0 and liquidez and liquidez > 0:
        result["Caja"] = liquidez / patrimonio

    cat1 = None
    for i in range(27, min(len(df), 300)):
        level = df.iloc[i, 0]
        if pd.isna(level):
            continue
        level = float(level)
        desc = str(df.iloc[i, 2]) if not pd.isna(df.iloc[i, 2]) else ""

        if level == 1.0:
            cat1 = desc
        elif level == 2.0 and cat1:
            pct = df.iloc[i, col_pct]
            if not pd.isna(pct) and isinstance(pct, (int, float)) and float(pct) > 0.001:
                # Normalize cat2 label (strip accents for map lookup)
                desc_norm = desc.replace("\u00f3", "o").replace("\u00e9", "e")
                label = COMP_PIE_MAP.get((cat1, desc)) or COMP_PIE_MAP.get((cat1, desc_norm))
                if label and label != "Caja":  # Caja ya agregada desde liquidez
                    result[label] = result.get(label, 0) + float(pct)

    # Normalizar para que sume 1
    total = sum(result.values())
    if total > 0.05:
        result = {k: v / total for k, v in result.items()}

    # Filtrar < 0.5% y ordenar por valor desc
    result = {k: v for k, v in result.items() if v >= 0.005}
    return dict(sorted(result.items(), key=lambda x: -x[1]))


def _get_top5(df: pd.DataFrame, col_pct: int) -> list:
    """
    Extrae los top 5 instrumentos por % de participación.
    Retorna lista de {ticker, pct, descripcion}.
    Excluye versiones CI (contado inmediato) y duplicados.
    """
    instrumentos = []
    seen = set()

    for i in range(27, min(len(df), 300)):
        nivel = df.iloc[i, 0]
        if pd.isna(nivel) or float(nivel) != 4.0:
            continue
        ticker = str(df.iloc[i, 1]) if not pd.isna(df.iloc[i, 1]) else ""
        pct    = df.iloc[i, col_pct]

        if not ticker or pd.isna(pct) or float(pct) < 0.001:
            continue

        # Limpiar sufijo CI para deduplicar
        ticker_base = ticker.replace(" CI", "").strip()
        if ticker_base in seen:
            continue
        seen.add(ticker_base)

        # Fecha de vencimiento desde col 3 (Av Life en años → approx date)
        desc = str(df.iloc[i, 2]) if not pd.isna(df.iloc[i, 2]) else ""

        # Intentar extraer fecha del nombre del bono (ej "Vto. 30 11 2026")
        fecha_vto = None
        m = re.search(r"[Vv]to[.\s]+(\d{1,2})[.\s/]+(\d{1,2})[.\s/]+(\d{4})", desc)
        if m:
            try:
                fecha_vto = f"{int(m.group(1)):02d}-{_MESES_ABREV[int(m.group(2))]}-{m.group(3)}"
            except Exception:
                pass

        instrumentos.append({
            "ticker":    ticker_base,
            "pct":       round(float(pct), 6),
            "fecha_vto": fecha_vto or "",
        })

    # Ordenar por % desc y tomar top 5
    instrumentos.sort(key=lambda x: -x["pct"])
    return instrumentos[:5]


_MESES_ABREV = {
    1:"ene.", 2:"feb.", 3:"mar.", 4:"abr.", 5:"may.", 6:"jun.",
    7:"jul.", 8:"ago.", 9:"sep.", 10:"oct.", 11:"nov.", 12:"dic."
}


def _get_fecha_datos(df: pd.DataFrame) -> str:
    """Extrae la fecha de datos del sheet (fila 25, col 2)."""
    try:
        v = df.iloc[25, 2]
        if pd.isna(v):
            return date.today().strftime("%d-%b-%Y")
        if hasattr(v, "strftime"):
            return v.strftime("%d-%b-%Y")
        return str(v)[:10]
    except Exception:
        return date.today().strftime("%d-%b-%Y")


def _get_liquidez(df: pd.DataFrame, col_metr: int) -> float:
    """Extrae el valor de Liquidez (row 11) del sheet."""
    try:
        val = df.iloc[10, col_metr]  # row 11 (0-indexed=10)
        return float(val) if not pd.isna(val) else 0.0
    except Exception:
        return 0.0


def _load_fondo(path: str, fondo: dict) -> dict:
    """Carga todos los datos de un fondo desde RF_Detalle."""
    df = _read_sheet(path, fondo["sheet_rf"])
    metr = _get_metrics(df, fondo["col_metr"])
    comp = _get_composition(df, fondo["col_pct"])   # para Tipo de Ajuste
    top5 = _get_top5(df, fondo["col_pct"])
    fecha = _get_fecha_datos(df)
    # Composicion general con detalle L1+L2
    patrimonio = metr.get("patrimonio") or 0
    liquidez   = _get_liquidez(df, fondo["col_metr"])
    comp_pie   = _get_composition_pie(df, fondo["col_pct"], fondo.get("col_monto"),
                                      patrimonio, liquidez)
    return {
        **fondo,
        "metrics": metr,
        "composition": comp,
        "comp_pie": comp_pie,
        "top5": top5,
        "fecha_datos": fecha,
    }


# =============================================================================
# FORMATEO
# =============================================================================

def _fmt_pn(v: float) -> str:
    if v is None: return "—"
    b = v / 1e9
    if b >= 1:
        return f"${b:,.2f}B".replace(",", "X").replace(".", ",").replace("X", ".")
    m = v / 1e6
    return f"${m:,.0f}M".replace(",", "X").replace(".", ",").replace("X", ".")

def _fmt_tir(v: float, is_cer: bool = False) -> str:
    if v is None: return "—"
    if is_cer:
        real = v - 0.20  # approx: TIR real ≈ TIR nominal - inflación esperada
        # Usar tir_neta directamente si disponible
        return f"CER+{real*100:.1f}%"
    return f"{v*100:.1f}%"

def _fmt_dur(v: float) -> str:
    if v is None: return "—"
    if v < 0.15:
        return f"{round(v * 365)} días"
    return f"{v:.2f} años"

def _fmt_pct(v: float) -> str:
    return f"{round(v * 100)}%"


# =============================================================================
# GENERADOR HTML — slides de fondos
# =============================================================================

_FONDO_CSS = """
/* ── Slides de fondos Delta — layout tipo PPTX ── */
.fd-slide-inner {
  flex: 1; overflow-y: auto;
  padding: 20px 32px 16px;
  font-family: 'Calibri', 'Segoe UI', Arial, sans-serif;
}
.fd-header {
  margin-bottom: 14px;
}
.fd-title {
  font-family: 'Calibri', 'Segoe UI', Arial, sans-serif;
  font-size: 26px; font-weight: 700; color: #0f2557; line-height: 1.1;
}
.fd-subtitle {
  font-size: 13px; color: #5f7080; margin-top: 2px;
  font-family: 'Calibri', 'Segoe UI', Arial, sans-serif;
}
/* Layout: 3 cols (KPI | composición+tipo | tenencias) + VCP abajo */
.fd-main {
  display: flex; flex-direction: column; gap: 10px;
  height: calc(100% - 64px); min-height: 380px;
}
.fd-top-row {
  display: grid;
  grid-template-columns: 140px 1fr 210px;
  gap: 10px; flex: 1; min-height: 0;
}
.fd-vcp-row {
  display: grid;
  grid-template-columns: 1fr 240px;
  gap: 10px; min-height: 185px; flex-shrink: 0;
}
/* Panel KPI */
.fd-kpi {
  background: #3f5060; border-radius: 6px;
  padding: 18px 14px;
  display: flex; flex-direction: column;
  justify-content: center; align-items: center;
  gap: 16px; color: white; text-align: center;
}
.fd-kpi-label {
  font-size: 10px; font-weight: 700;
  text-transform: uppercase; letter-spacing: .8px;
  color: rgba(255,255,255,.65); margin-bottom: 3px;
}
.fd-kpi-value {
  font-size: 21px; font-weight: 700; color: white; line-height: 1.1;
}
.fd-kpi-date { font-size: 11px; color: rgba(255,255,255,.45); margin-top: 2px; }
.fd-charts-col {
  display: flex; flex-direction: column; gap: 10px; min-height: 0;
}
/* Quadrante de contenido */
.fd-quadrant {
  background: white; border: 1px solid #d0d9e4;
  border-radius: 6px; overflow: hidden;
  display: flex; flex-direction: column;
}
.fd-q-title {
  background: #0f2557; color: white;
  font-family: 'Calibri','Segoe UI',Arial,sans-serif;
  font-size: 10px; font-weight: 700;
  text-transform: uppercase; letter-spacing: .5px;
  padding: 5px 10px; flex-shrink: 0;
}
.fd-q-body {
  flex: 1; overflow: hidden;
  display: flex; align-items: center; justify-content: center;
  padding: 8px;
}
/* Top 5 table */
.fd-top5 {
  width: 100%; border-collapse: collapse;
  font-family: 'Calibri','Segoe UI',Arial,sans-serif;
  font-size: 12px;
}
.fd-top5 tr { border-bottom: 1px solid #edf0f3; }
.fd-top5 td { padding: 4px 6px; }
.fd-top5 .tk { font-weight: 700; color: #0f2557; }
.fd-top5 .pct {
  background: #5f7080; color: white; border-radius: 3px;
  padding: 2px 7px; font-weight: 700; font-size: 11px;
  text-align: center; white-space: nowrap;
}
.fd-top5 .fecha { color: #8090a0; font-size: 11px; }
/* VCP */
.fd-vcp-form {
  width: 100%;
  font-family: 'Calibri','Segoe UI',Arial,sans-serif;
  font-size: 11px;
}
.fd-vcp-note {
  font-size: 9px; color: #8090a0; margin-bottom: 4px;
}
.fd-vcp-grid {
  display: grid; grid-template-columns: 90px 1fr 1fr;
  gap: 3px 6px; align-items: center;
}
.fd-vcp-grid label { color: #5f7080; font-size: 10.5px; }
.fd-vcp-grid .hdr { font-size: 9px; font-weight: 700; color: #7a90a4; text-transform: uppercase; }
.fd-vcp-grid input {
  width: 100%; border: 1px solid #d0d9e4; border-radius: 3px;
  padding: 2px 5px; font-size: 11px; text-align: center;
  font-family: 'Calibri','Segoe UI',Arial,sans-serif;
  box-sizing: border-box;
}
"""


def _pie_chart_js(canvas_id: str, labels: list, values: list, colors: list) -> str:
    """Pie chart con etiquetas externas y lineas conectoras, estilo PPTX."""
    data_js   = json.dumps([round(v * 100, 2) for v in values])
    labels_js = json.dumps(labels)
    colors_js = json.dumps(colors)

    return f"""
<script>
(function(){{
  function drawPie_{canvas_id}(){{
    var cv = document.getElementById('{canvas_id}');
    if (!cv) return;
    var W = cv.parentElement.clientWidth;
    var H = cv.parentElement.clientHeight;
    if (!W || W < 10 || !H || H < 10) {{ setTimeout(drawPie_{canvas_id}, 200); return; }}
    var dpr = window.devicePixelRatio || 1;
    cv.width = W * dpr; cv.height = H * dpr;
    cv.style.width = W + 'px'; cv.style.height = H + 'px';
    var ctx = cv.getContext('2d');
    ctx.scale(dpr, dpr);
    ctx.clearRect(0, 0, W, H);

    var allData   = {data_js};
    var allLabels = {labels_js};
    var allColors = {colors_js};

    // Filtrar slices < 1%
    var data = [], labels = [], colors = [];
    var tot0 = allData.reduce(function(a,b){{return a+b;}},0);
    for (var k=0;k<allData.length;k++) {{
      if (allData[k]/tot0 >= 0.01) {{
        data.push(allData[k]);
        labels.push(allLabels[k]);
        colors.push(allColors[k]);
      }}
    }}
    var total = data.reduce(function(a,b){{return a+b;}},0);
    if (total <= 0) return;

    // Margen para etiquetas: izquierda y derecha
    var margin = Math.min(W * 0.32, 90);
    var cx = W / 2;
    var cy = H / 2;
    var R  = Math.min((W - margin*2) / 2, H / 2) * 0.82;
    R = Math.max(R, 30);

    // Calcular angulos
    var slices = [];
    var angle = -Math.PI / 2;
    for (var i=0;i<data.length;i++) {{
      var s = (data[i]/total) * 2 * Math.PI;
      slices.push({{ start:angle, end:angle+s, pct:data[i]/total, mid:angle+s/2 }});
      angle += s;
    }}

    // Dibujar slices
    for (var i=0;i<slices.length;i++) {{
      var sl = slices[i];
      ctx.beginPath();
      ctx.moveTo(cx,cy);
      ctx.arc(cx,cy,R,sl.start,sl.end);
      ctx.closePath();
      ctx.fillStyle = colors[i];
      ctx.fill();
      ctx.strokeStyle='white';
      ctx.lineWidth=1.5;
      ctx.stroke();
    }}

    // Etiquetas externas con lineas conectoras
    // Separar left/right por lado del canvas
    var leftItems  = [];
    var rightItems = [];
    for (var i=0;i<slices.length;i++) {{
      var mid = slices[i].mid;
      var onRight = (Math.cos(mid) >= 0);
      (onRight ? rightItems : leftItems).push(i);
    }}

    // Funcion para dibujar un lado
    function drawSide(indices, isRight) {{
      if (indices.length === 0) return;
      var side = isRight ? 1 : -1;

      // Anchor points en el borde del pie
      var anchors = indices.map(function(i) {{
        var mid = slices[i].mid;
        return {{
          idx: i,
          ax: cx + Math.cos(mid) * R,
          ay: cy + Math.sin(mid) * R,
          mid: mid,
          pct: slices[i].pct
        }};
      }});

      // Punto de elbow: un paso afuera del pie
      anchors.forEach(function(a) {{
        a.ex = cx + Math.cos(a.mid) * (R + 10);
        a.ey = cy + Math.sin(a.mid) * (R + 10);
      }});

      // Y final de la etiqueta: distribuido verticalmente sin solapamiento
      var lx = isRight ? cx + R + 22 : cx - R - 22;
      var lineH = 13;
      var totalH = (anchors.length - 1) * lineH;
      var startY = cy - totalH / 2;

      // Ordenar por angulo (arriba a abajo por y)
      anchors.sort(function(a,b){{ return a.ey - b.ey; }});
      anchors.forEach(function(a, j) {{
        a.ly = startY + j * lineH;
      }});

      // Dibujar cada etiqueta
      anchors.forEach(function(a) {{
        var i = a.idx;
        ctx.strokeStyle = colors[i];
        ctx.lineWidth = 0.9;
        ctx.beginPath();
        ctx.moveTo(a.ax, a.ay);
        ctx.lineTo(a.ex, a.ey);
        ctx.lineTo(lx + (isRight ? -4 : 4), a.ly);
        ctx.stroke();

        // Punto en el quiebre
        ctx.beginPath();
        ctx.arc(a.ax, a.ay, 1.5, 0, 2*Math.PI);
        ctx.fillStyle = colors[i];
        ctx.fill();

        // Texto: "Categoria; X%"
        var pctStr = (a.pct * 100).toFixed(a.pct < 0.05 ? 1 : 0) + '%';
        var txt = labels[i] + '; ' + pctStr;
        ctx.font = '9px Barlow,Calibri,Arial';
        ctx.fillStyle = '#1a1a2e';
        ctx.textAlign = isRight ? 'left' : 'right';
        ctx.textBaseline = 'middle';
        ctx.fillText(txt, lx, a.ly);
      }});
    }}

    drawSide(rightItems, true);
    drawSide(leftItems, false);
  }}
  window.addEventListener('load', drawPie_{canvas_id});
  window.addEventListener('resize', drawPie_{canvas_id});
  window['drawPie_{canvas_id}'] = drawPie_{canvas_id};
  // Disparar al navegar al slide
  (function patchNav_{canvas_id}() {{
    var _p = window.goToId;
    if (typeof _p !== 'function') {{ setTimeout(patchNav_{canvas_id}, 100); return; }}
    var _cid = '{canvas_id}';
    window.goToId = function(id) {{
      _p(id);
      // Re-draw cuando se activa este fd-slide
      if (id && _cid && id.indexOf(_cid.replace('pie_','')) !== -1) {{
        setTimeout(drawPie_{canvas_id}, 100);
        var bf = window['drawBar_bar_' + _cid.replace('pie_','')];
        if (typeof bf === 'function') setTimeout(bf, 100);
      }}
    }};
  }})();
}})();
</script>
"""


def _bar_chart_js(canvas_id: str, labels: list, values: list, color: str) -> str:
    """Barra horizontal para tipo de ajuste."""
    data_js   = json.dumps([round(v * 100, 2) for v in values])
    labels_js = json.dumps(labels)

    return f"""
<script>
(function(){{
  function drawBar_{canvas_id}(){{
    var cv = document.getElementById('{canvas_id}');
    if (!cv) return;
    var W = cv.parentElement.clientWidth;
    var H = cv.parentElement.clientHeight;
    if (!W || W < 10) {{ setTimeout(drawBar_{canvas_id}, 200); return; }}
    var dpr = window.devicePixelRatio || 1;
    cv.width = W * dpr; cv.height = H * dpr;
    cv.style.width = W + 'px'; cv.style.height = H + 'px';
    var ctx = cv.getContext('2d');
    ctx.scale(dpr, dpr);
    ctx.fillStyle = 'white'; ctx.fillRect(0, 0, W, H);

    var data   = {data_js};
    var labels = {labels_js};
    var n = labels.length;
    var ML = 78, MR = 48, MT = 6, MB = 6;
    var PW = W - ML - MR, PH = H - MT - MB;
    var rowH = PH / n;
    var maxVal = Math.max.apply(null, data) || 1;

    for (var i = 0; i < n; i++) {{
      var y = MT + i * rowH;
      var bw = (data[i] / 100) * PW;
      // bg
      ctx.fillStyle = '#f0f4f8';
      ctx.fillRect(ML, y + rowH*0.15, PW, rowH * 0.7);
      // bar
      ctx.fillStyle = '{color}';
      ctx.fillRect(ML, y + rowH*0.15, bw, rowH * 0.7);
      // label izq
      ctx.font = '9.5px Barlow,Calibri,Arial';
      ctx.fillStyle = '#333';
      ctx.textAlign = 'right';
      ctx.textBaseline = 'middle';
      ctx.fillText(labels[i], ML - 4, y + rowH / 2);
      // valor der
      ctx.textAlign = 'left';
      ctx.fillStyle = data[i] > 5 ? 'white' : '#333';
      if (data[i] > 8) {{
        ctx.fillText(data[i].toFixed(1) + '%', ML + bw - 28, y + rowH / 2);
      }} else {{
        ctx.fillStyle = '#333';
        ctx.fillText(data[i].toFixed(1) + '%', ML + bw + 3, y + rowH / 2);
      }}
    }}
  }}
  window.addEventListener('load', drawBar_{canvas_id});
  window.addEventListener('resize', drawBar_{canvas_id});
}})();
</script>
"""


def _vcp_chart_js(fondo_id: str) -> str:
    """JS para el gráfico de variaciones VCP — lee inputs del formulario."""
    return f"""
<script>
(function(){{
  function drawVcp_{fondo_id}(){{
    var cv = document.getElementById('vcp-cv-{fondo_id}');
    if (!cv) return;
    var W = cv.parentElement.clientWidth;
    if (!W || W < 10) {{ setTimeout(drawVcp_{fondo_id}, 80); return; }}
    var H = cv.parentElement.clientHeight > 40 ? cv.parentElement.clientHeight - 4 : 160;
    var dpr = window.devicePixelRatio || 1;
    cv.width = W * dpr; cv.height = H * dpr;
    cv.style.width = W + 'px'; cv.style.height = H + 'px';
    var ctx = cv.getContext('2d');
    ctx.scale(dpr, dpr);
    ctx.fillStyle = 'white'; ctx.fillRect(0, 0, W, H);

    // Leer valores del formulario
    function v(id){{ var el=document.getElementById(id); return el?parseFloat(el.value)||0:0; }}
    var entities = ['fondo','benchmark','peers','industria'];
    var labels   = [
      document.getElementById('vcp-lbl-{fondo_id}') ?
        document.getElementById('vcp-lbl-{fondo_id}').value || '{fondo_id}' : '{fondo_id}',
      'Benchmark','Peers','Industria'
    ];
    var s1 = entities.map(function(e){{ return v('vcp-{fondo_id}-s1-'+e); }});
    var s30= entities.map(function(e){{ return v('vcp-{fondo_id}-s30-'+e); }});

    var all = s1.concat(s30);
    var maxV = Math.max.apply(null, all.map(Math.abs)) || 0.01;
    var n = 4;
    var ML = 8, MR = 8, MT = 18, MB = 22;
    var PW = W - ML - MR, PH = H - MT - MB;
    var grpW = PW / n;
    var barW = grpW * 0.33;

    // Eje 0
    var zero = MT + PH * maxV / (2 * maxV);
    ctx.strokeStyle = '#d8e2ed'; ctx.lineWidth = 0.8;
    ctx.beginPath(); ctx.moveTo(ML, zero); ctx.lineTo(ML+PW, zero); ctx.stroke();

    var colors1  = '#1e6fba';
    var colors30 = '#1a7a46';

    for (var i = 0; i < n; i++) {{
      var gx = ML + i * grpW + grpW * 0.1;

      // Barra semana
      var h1 = s1[i] / maxV * PH / 2;
      ctx.fillStyle = colors1;
      if (h1 >= 0) ctx.fillRect(gx, zero - h1, barW, h1);
      else ctx.fillRect(gx, zero, barW, -h1);

      // Barra 30d
      var h30 = s30[i] / maxV * PH / 2;
      ctx.fillStyle = colors30;
      if (h30 >= 0) ctx.fillRect(gx + barW + 2, zero - h30, barW, h30);
      else ctx.fillRect(gx + barW + 2, zero, barW, -h30);

      // Label valor semana
      ctx.font = 'bold 10px Barlow,Calibri,Arial';
      ctx.fillStyle = colors1; ctx.textAlign = 'center';
      var vsy = h1 >= 0 ? zero - h1 - 3 : zero + (-h1) + 10;
      ctx.fillText(s1[i].toFixed(1)+'%', gx + barW/2, vsy);

      // Label valor 30d
      ctx.fillStyle = colors30;
      var v30y = h30 >= 0 ? zero - h30 - 3 : zero + (-h30) + 10;
      ctx.fillText(s30[i].toFixed(1)+'%', gx + barW + 2 + barW/2, v30y);

      // Label entidad
      ctx.fillStyle = '#333'; ctx.font = '10px Barlow,Calibri,Arial';
      ctx.fillText(labels[i], gx + barW + 2, H - 6);
    }}

    // Leyenda
    ctx.fillStyle = colors1;  ctx.fillRect(ML, 5, 10, 8);
    ctx.fillStyle = '#333'; ctx.font = '10px Barlow,Calibri,Arial'; ctx.textAlign='left';
    ctx.fillText('Últ. semana', ML+11, 12);
    ctx.fillStyle = colors30; ctx.fillRect(ML+80, 5, 8, 7);
    ctx.fillText('Últ. 30d', ML+92, 12);
  }}
  window.addEventListener('load', drawVcp_{fondo_id});
  // Redibujar cuando cambian inputs
  document.addEventListener('input', function(e){{
    if (e.target && e.target.id && e.target.id.indexOf('vcp-{fondo_id}') === 0)
      drawVcp_{fondo_id}();
  }});
  window.addEventListener('resize', drawVcp_{fondo_id});
  window['drawVcp_{fondo_id}'] = drawVcp_{fondo_id};
}})();
</script>
"""


def _build_fondo_slide(fondo_data: dict, slide_num: int, total_slides: int) -> str:
    fd   = fondo_data
    metr = fd["metrics"]
    comp = fd["composition"]
    top5 = fd["top5"]
    fid  = fd["id"]

    pn_str  = _fmt_pn(metr.get("patrimonio"))
    dur_str = _fmt_dur(metr.get("duration"))
    is_cer  = (fd["sheet_rf"] == "CER")
    tir_neta = metr.get("tir_neta") or metr.get("tir")
    tir_str = (f"CER+{tir_neta*100:.1f}%" if is_cer else f"{tir_neta*100:.1f}%") if tir_neta else "—"
    fecha_str = fd.get("fecha_datos", "—")

    # Composicion general: usar comp_pie (L1+L2 detallado)
    comp_pie = fd.get("comp_pie", comp)  # fallback a comp si no hay comp_pie
    pie_labels_raw = list(comp_pie.keys())
    pie_values_raw = list(comp_pie.values())
    pie_colors = [COMP_PIE_COLORS.get(k, "#888") for k in pie_labels_raw]
    pie_js = _pie_chart_js(f"pie_{fid}", pie_labels_raw, pie_values_raw, pie_colors)
    # Tipo de ajuste: usar comp (AJUSTE_ORDER) ordenado de mayor a menor
    bar_labels = sorted(comp.keys(), key=lambda k: -comp[k])
    bar_values = [comp[k] for k in bar_labels]
    bar_js = _bar_chart_js(f"bar_{fid}", bar_labels, bar_values, "#1e6fba")
    vcp_js = _vcp_chart_js(fid)

    top5_rows = "".join(
        f'<tr><td class="tk">{t["ticker"]}</td>'
        f'<td><span class="pct">{round(t["pct"]*100)}%</span></td>'
        f'<td class="fecha">{t.get("fecha_vto","")}</td></tr>'
        for t in top5
    )

    entities = [("fondo", fd["peers_label"]),("benchmark","Benchmark"),
                ("peers","Peers"),("industria","Industria")]
    vcp_rows = "\n".join(
        f'<label>{lbl}</label>'
        f'<input type="number" step="0.01" id="vcp-{fid}-s1-{eid}" placeholder="0.0" oninput="drawVcp_{fid}()">'
        f'<input type="number" step="0.01" id="vcp-{fid}-s30-{eid}" placeholder="0.0" oninput="drawVcp_{fid}()">'
        for eid, lbl in entities
    )

    fondo_ids = [f["id"] for f in FONDOS]
    prev_id = f"fd-slide-{fondo_ids[slide_num-2]}" if slide_num > 1 else "slide-6"
    next_id = f"fd-slide-{fondo_ids[slide_num]}"   if slide_num < total_slides else "slide-7"
    pct_prog = round(slide_num / total_slides * 100, 1)
    kpi_tir = (f'<div class="fd-kpi-item"><div class="fd-kpi-label">TIR*</div>'
               f'<div class="fd-kpi-value">{tir_str}</div></div>') if fd.get("tiene_tir") else ""

    return (
        f'\n<div class="slide-chapter fd-slide" id="fd-slide-{fid}">\n'
        f'  <div class="progress-bar"><div class="progress-fill" style="width:{pct_prog}%"></div></div>\n'
        f'  <div class="slide-chapter-inner fd-slide-inner">\n'
        f'    <div class="fd-header">\n'
        f'      <div class="fd-title">{fd["nombre"]}</div>\n'
        f'      <div class="fd-subtitle">Performance y posicionamiento</div>\n'
        f'    </div>\n'
        f'    <div class="fd-main">\n'
        f'      <div class="fd-top-row">\n'
        f'        <div class="fd-kpi">\n'
        f'          <div><div class="fd-kpi-label">Patrimonio</div>'
        f'<div class="fd-kpi-value">{pn_str}</div></div>\n'
        f'          {kpi_tir}\n'
        f'          <div><div class="fd-kpi-label">Duration*</div>'
        f'<div class="fd-kpi-value">{dur_str}</div></div>\n'
        f'          <div><div class="fd-kpi-label">Datos al</div>'
        f'<div class="fd-kpi-date">{fecha_str}</div></div>\n'
        f'        </div>\n'
        f'        <div class="fd-charts-col">\n'
        f'          <div class="fd-quadrant" style="flex:1"><div class="fd-q-title">Composición General</div>'
        f'<div class="fd-q-body"><canvas id="pie_{fid}" style="width:100%;height:100%;display:block"></canvas></div></div>\n'
        f'          <div class="fd-quadrant" style="flex:1"><div class="fd-q-title">Tipo de Ajuste</div>'
        f'<div class="fd-q-body"><canvas id="bar_{fid}" style="width:100%;height:100%;display:block"></canvas></div></div>\n'
        f'        </div>\n'
        f'        <div class="fd-quadrant"><div class="fd-q-title">Principales Tenencias</div>'
        f'<div class="fd-q-body" style="align-items:flex-start;padding:12px 14px">'
        f'<table class="fd-top5" style="font-size:13px"><tbody>{top5_rows}</tbody></table></div></div>\n'
        f'      </div>\n'
        f'      <div class="fd-vcp-row">\n'
        f'        <div class="fd-quadrant"><div class="fd-q-title">Variaciones VCP — {fd["tipo_vcp"]}</div>'
        f'<div class="fd-q-body" style="padding:8px 12px;flex-direction:column;align-items:stretch">'
        f'<canvas id="vcp_cv_{fid}" style="width:100%;flex:1"></canvas></div></div>\n'
        f'        <div class="fd-quadrant"><div class="fd-q-title">Ingresá valores</div>'
        f'<div class="fd-q-body" style="flex-direction:column;align-items:stretch;padding:8px 10px">'
        f'<div class="fd-vcp-note" style="margin-bottom:8px">✏ Valores en % (ej: 0.9 = 0.9%)</div>'
        f'<div class="fd-vcp-grid">'
        f'<span class="hdr"></span><span class="hdr">Últ. semana</span><span class="hdr">Últ. 30d</span>'
        f'<input type="hidden" id="vcp-lbl-{fid}" value="{fd["peers_label"]}">\n'
        f'{vcp_rows}\n</div></div></div>\n'
        f'      </div>\n'
        f'    </div>\n  </div>\n'
        f'  <div class="slide-footer">'
        f'<div class="slide-footer-left">Delta Asset Management · {fecha_str} · Uso interno</div>'
        f'<div class="slide-nav-btns">'
        f'<button class="slide-nav-btn" onclick="goToId(\'{prev_id}\')">← Anterior</button>'
        f'<button class="slide-nav-btn" onclick="goToId(\'{next_id}\')">Siguiente →</button>'
        f'</div><div class="slide-footer-right">{slide_num} / {total_slides}</div></div>\n'
        f'</div>\n{pie_js}\n{bar_js}\n{vcp_js}\n'
    )



def generar_fondos_nav_items(fondos: list = None) -> str:
    """
    Genera los <div class="nav-item"> para el sidebar del reporte HTML.
    Se insertan después del nav-item 07 (Snapshot).
    Cada item usa goToId() para navegar a los fd-slides.
    """
    fondos = fondos or FONDOS
    n = len(fondos)
    items = []
    for i, fd in enumerate(fondos, 1):
        num = f"{7 + i:02d}"
        items.append(f"""    <div class="nav-item fd-nav-item" onclick="goToId('fd-slide-{fd['id']}')" id="fd-nav-{fd['id']}">
      <span class="nav-num">{num}</span>
      <span class="nav-label">{fd['nombre']}<br>
        <span style="font-size:10px;opacity:.6">Composición · Top 5</span>
      </span>
    </div>""")
    # Nav item para slide 12 - Inputs VCP
    num_inp = f"{7 + len(fondos) + 1:02d}"
    items.append(
        f'''    <div class="nav-item fd-nav-item" onclick="goToId('fd-slide-inputs')" id="fd-nav-inputs">\n'''
        f'''      <span class="nav-num">{num_inp}</span>\n'''
        '''      <span class="nav-label">Inputs VCP<br>\n'''
        '''        <span style="font-size:10px;opacity:.6">Variaciones - Todos los fondos</span>\n'''
        '''      </span>\n'''
        '''    </div>'''
    )
    return "\n".join(items)


def generar_fondos_html(
    rf_detalle_path: str,
    fondos: list = None,
) -> str:
    """
    Genera el HTML de las slides de fondos para embeber en el reporte.
    Retorna el HTML completo (CSS + slides + JS).
    """
    fondos = fondos or FONDOS
    slides_html = ""
    n = len(fondos)

    for i, fondo_cfg in enumerate(fondos, 1):
        try:
            fd = _load_fondo(rf_detalle_path, fondo_cfg)
            slides_html += _build_fondo_slide(fd, i, n)
            print(f"[generar_fondos] ✓ {fondo_cfg['nombre']}: "
                  f"PN={fd['metrics'].get('patrimonio','?'):.2e} "
                  f"comp={list(fd['composition'].keys())[:3]} "
                  f"top5={[t['ticker'] for t in fd['top5']]}")
        except Exception as e:
            import traceback
            print(f"[generar_fondos] ✗ Error en {fondo_cfg['nombre']}: {e}")
            traceback.print_exc()
            continue

    # Script de navegación goToId() — compatibiliza con goTo() del reporte
    goto_js = """
<script>
// goToId — wrapper de la navegación unificada
// La función goToId real está definida en el HTML por _unify_navigation()
// Este script solo se incluye como fallback por si acaso
if (typeof goToId === 'undefined') {
  function goToId(targetId) {
    var m = targetId.match(/^slide-([0-9]+)$/);
    if (m && typeof goTo === 'function') { goTo(parseInt(m[1])); return; }
    document.querySelectorAll('.slide-chapter, .slide-cover').forEach(function(s) {
      s.classList.remove('active');
    });
    var target = document.getElementById(targetId);
    if (target) {
      target.classList.add('active');
      setTimeout(function() {
        window.dispatchEvent(new Event('resize'));
        var active = document.querySelector('.fd-slide.active');
        if (active) {
          var fid = active.id.replace('fd-slide-','');
          if (typeof window['drawVcp_'+fid]==='function') window['drawVcp_'+fid]();
        }
      }, 100);
    }
    document.querySelectorAll('.nav-item').forEach(function(n) { n.classList.remove('active'); });
    var navId = targetId.replace('fd-slide-', 'fd-nav-').replace(/^slide-/, 'nav-');
    var navEl = document.getElementById(navId);
    if (navEl) navEl.classList.add('active');
  }
}
</script>
"""
    return f"<style>{_FONDO_CSS}</style>\n{slides_html}\n{goto_js}"


# =============================================================================
# API PÚBLICA — PPTX
# =============================================================================


def generar_inputs_slide(fondos=None):
    """Slide 12: inputs VCP para todos los fondos. Soporta pre-carga desde CSV."""
    fondos = fondos or FONDOS
    lbl_s1  = "Ult. semana (%)"
    lbl_s30 = "Ult. 30d (%)"

    rows_html = ""
    for fd in fondos:
        fid = fd["id"]
        entities = [
            ("fondo",     fd.get("peers_label", fd["nombre"])),
            ("benchmark", "Benchmark"),
            ("peers",     "Peers"),
            ("industria", "Industria"),
        ]
        tbl_rows = ""
        inp_style = "width:90px;border:1px solid #c0cdd8;border-radius:4px;padding:4px 8px;font-size:13px;text-align:center"
        for eid, lbl in entities:
            tbl_rows += "".join([
                "<tr style='border-bottom:1px solid #edf0f3'>",
                f"<td style='padding:6px 12px;font-weight:600;color:#333;font-size:13px;min-width:130px'>{lbl}</td>",
                "<td style='padding:4px 8px;text-align:center'>",
                f"<input type='number' step='0.01' placeholder='0.0' id='vcp-{fid}-s1-{eid}'",
                f" oninput='if(typeof drawVcp_{fid}==\"function\") drawVcp_{fid}()'",
                f" style='{inp_style}'></td>",
                "<td style='padding:4px 8px;text-align:center'>",
                f"<input type='number' step='0.01' placeholder='0.0' id='vcp-{fid}-s30-{eid}'",
                f" oninput='if(typeof drawVcp_{fid}==\"function\") drawVcp_{fid}()'",
                f" style='{inp_style}'></td>",
                "</tr>",
            ])

        tipo = fd.get("tipo_vcp", "")
        th_style = "padding:6px 12px;text-align:center;font-size:10px;color:#666;font-weight:700;text-transform:uppercase;border-bottom:1px solid #d0d9e4"
        rows_html += "".join([
            "<div style='background:white;border:1px solid #d0d9e4;border-radius:6px;overflow:hidden;margin-bottom:16px'>",
            "<div style='background:#0f2557;color:white;font-size:11px;font-weight:700;text-transform:uppercase;letter-spacing:.5px;padding:7px 14px'>",
            f"{fd['nombre']} <span style='font-weight:400;opacity:.7;font-size:10px;text-transform:none'>{tipo}</span>",
            "</div>",
            "<table style='width:100%;border-collapse:collapse;font-family:Calibri,Arial,sans-serif'>",
            "<thead><tr style='background:#f5f7fa'>",
            "<th style='padding:6px 12px;text-align:left;font-size:10px;color:#666;font-weight:700;text-transform:uppercase;border-bottom:1px solid #d0d9e4;min-width:130px'></th>",
            f"<th style='{th_style}'>{lbl_s1}</th>",
            f"<th style='{th_style}'>{lbl_s30}</th>",
            "</tr></thead>",
            f"<tbody>{tbl_rows}</tbody>",
            "</table></div>",
        ])

    slide_js = "\n".join([
        "<script>",
        "(function(){",
        "  function sizeInp(){",
        "    var s=document.getElementById('fd-slide-inputs');",
        "    var sc=document.getElementById('inp-scroll');",
        "    if(!s||!sc||!s.classList.contains('active'))return;",
        "    s.style.flexDirection='column';",
        "    var h=s.getBoundingClientRect().height||window.innerHeight;",
        "    var pb=s.querySelector('.progress-bar');",
        "    var ft=s.querySelector('.slide-footer');",
        "    var av=h-(pb?pb.offsetHeight:3)-(ft?ft.offsetHeight:36)-2;",
        "    if(av>100){sc.style.height=av+'px';sc.style.flex='none';}",
        "  }",
        "  function patch(){",
        "    var p=window.goToId;",
        "    if(typeof p!=='function'){setTimeout(patch,100);return;}",
        "    window.goToId=function(id){p(id);if(id==='fd-slide-inputs'){setTimeout(sizeInp,50);setTimeout(sizeInp,300);}};",
        "  }",
        "  patch();",
        "  window.addEventListener('resize',function(){",
        "    var s=document.getElementById('fd-slide-inputs');",
        "    if(s&&s.classList.contains('active'))sizeInp();",
        "  });",
        "})();",
        "</script>",
    ])

    css = (
        "<style>"
        "#fd-slide-inputs{display:none;height:100%;background:white;flex-direction:column;overflow:hidden}"
        "#fd-slide-inputs.active{display:flex!important;flex-direction:column!important}"
        "</style>"
    )

    btn = "goToId(\"slide-7\")"
    return "\n".join([
        css,
        '<div class="slide-chapter fd-slide" id="fd-slide-inputs">',
        '  <div class="progress-bar"><div class="progress-fill" style="width:100%"></div></div>',
        '  <div id="inp-scroll" style="overflow-y:auto;padding:24px 36px 16px;font-family:Calibri,Arial,sans-serif">',
        '    <div style="font-size:24px;font-weight:700;color:#0f2557;margin-bottom:6px">Inputs - Variaciones VCP</div>',
        '    <div style="font-size:13px;color:#5f7080;margin-bottom:4px">Completa los valores. Los graficos se actualizan automaticamente.</div>',
        '    <div style="font-size:11px;color:#7a90a4;font-style:italic;margin-bottom:18px">Valores en % directo (ej: 0.9 = 0.9%)</div>',
        '    <div style="display:grid;grid-template-columns:repeat(auto-fill,minmax(340px,1fr));gap:18px">',
        f'      {rows_html}',
        '    </div>',
        '  </div>',
        '  <div class="slide-footer">',
        '    <div class="slide-footer-left">Delta Asset Management - Inputs VCP</div>',
        f'    <div class="slide-nav-btns"><button class="slide-nav-btn" onclick="{btn}">Volver al reporte</button></div>',
        '    <div class="slide-footer-right">12</div>',
        '  </div>',
        '</div>',
        slide_js,
    ])


def generar_fondos_pptx(
    rf_detalle_path: str,
    output_path: str,
    template_path: str = None,
    fondos: list = None,
) -> str:
    """
    Genera una presentación PPTX con las slides de fondos.
    Si se provee template_path, la usa como base (mantiene el diseño).
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        import io
    except ImportError:
        print("[generar_fondos] python-pptx no instalado. Instalá con: pip install python-pptx")
        return None

    fondos = fondos or FONDOS

    if template_path and Path(template_path).exists():
        prs = Presentation(template_path)
        # Eliminar slides existentes para reemplazar con datos frescos
        NS = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
        xml_slides = prs.slides._sldIdLst
        for sld_id in list(xml_slides):
            r_id = sld_id.get(f'{{{NS}}}id')
            if r_id:
                try:
                    prs.part.drop_rel(r_id)
                except Exception:
                    pass
            xml_slides.remove(sld_id)
    else:
        prs = Presentation()
        prs.slide_width  = Inches(13.33)
        prs.slide_height = Inches(7.5)

    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    # Colores Delta
    NAVY   = RGBColor(0x0f, 0x25, 0x57)
    GRAY   = RGBColor(0x5f, 0x70, 0x80)
    WHITE  = RGBColor(0xff, 0xff, 0xff)
    BLUE   = RGBColor(0x1e, 0x6f, 0xba)

    # Usar layout en blanco (último disponible) o el primero
    n_layouts = len(prs.slide_layouts)
    slide_layout = prs.slide_layouts[min(6, n_layouts - 1)]

    for fondo_cfg in fondos:
        try:
            fd = _load_fondo(rf_detalle_path, fondo_cfg)
        except Exception as e:
            print(f"[generar_fondos] PPTX error {fondo_cfg['nombre']}: {e}")
            continue

        slide = prs.slides.add_slide(slide_layout)
        metr  = fd["metrics"]
        comp  = fd["composition"]
        top5  = fd["top5"]

        # ── Título ────────────────────────────────────────────────────
        txb = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(9), Inches(0.55))
        tf  = txb.text_frame
        p   = tf.paragraphs[0]
        run = p.add_run()
        run.text = fd["nombre"]
        run.font.bold   = True
        run.font.size   = Pt(28)
        run.font.color.rgb = NAVY
        run.font.name   = "Barlow"

        txb2 = slide.shapes.add_textbox(Inches(0.4), Inches(0.72), Inches(9), Inches(0.3))
        p2   = txb2.text_frame.paragraphs[0]
        r2   = p2.add_run()
        r2.text = "Performance y posicionamiento"
        r2.font.size = Pt(13)
        r2.font.color.rgb = GRAY
        r2.font.name = "Calibri"

        # ── Panel KPI (rect gris) ─────────────────────────────────────
        kpi_box = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            Inches(0.25), Inches(1.1), Inches(2.1), Inches(5.8)
        )
        kpi_box.fill.solid()
        kpi_box.fill.fore_color.rgb = GRAY
        kpi_box.line.fill.background()

        def _add_kpi(y_inch, label, value):
            lbl = slide.shapes.add_textbox(Inches(0.35), Inches(y_inch), Inches(1.9), Inches(0.22))
            lp  = lbl.text_frame.paragraphs[0]
            lr  = lp.add_run()
            lr.text = label
            lr.font.size = Pt(8); lr.font.bold = True
            lr.font.color.rgb = RGBColor(0xd0, 0xd8, 0xe0)
            lr.font.name = "Barlow"

            val = slide.shapes.add_textbox(Inches(0.35), Inches(y_inch + 0.22), Inches(1.9), Inches(0.35))
            vp  = val.text_frame.paragraphs[0]
            vr  = vp.add_run()
            vr.text = value
            vr.font.size = Pt(16); vr.font.bold = True
            vr.font.color.rgb = WHITE
            vr.font.name = "Barlow"

        is_cer = fd["sheet_rf"] == "CER"
        tir_neta = metr.get("tir_neta") or metr.get("tir")
        if is_cer:
            tir_str = f"CER+{tir_neta*100:.1f}%" if tir_neta else "—"
        else:
            tir_str = f"{tir_neta*100:.1f}%" if tir_neta else "—"

        _add_kpi(1.3,  "Patrimonio",  _fmt_pn(metr.get("patrimonio")))
        if fd.get("tiene_tir"):
            _add_kpi(2.15, "TIR*", tir_str)
        _add_kpi(3.0,  "Duration*", _fmt_dur(metr.get("duration")))
        _add_kpi(3.85, "Datos al",  fd.get("fecha_datos", "—"))

        # ── Top 5 Tenencias ───────────────────────────────────────────
        t5_x = Inches(2.55); t5_y = Inches(4.1)
        hdr = slide.shapes.add_textbox(t5_x, Inches(3.85), Inches(4.3), Inches(0.24))
        hp  = hdr.text_frame.paragraphs[0]
        hr  = hp.add_run()
        hr.text = "Principales Tenencias"
        hr.font.size = Pt(9); hr.font.bold = True
        hr.font.color.rgb = NAVY; hr.font.name = "Barlow"

        for j, t in enumerate(top5):
            row_y = t5_y.inches + j * 0.28
            # Ticker
            tb_t = slide.shapes.add_textbox(Inches(2.55), Inches(row_y), Inches(1.2), Inches(0.26))
            pp   = tb_t.text_frame.paragraphs[0]
            rr   = pp.add_run()
            rr.text = t["ticker"]; rr.font.size = Pt(10.5)
            rr.font.bold = True; rr.font.color.rgb = NAVY; rr.font.name = "Calibri"
            # Pct badge
            pct_box = slide.shapes.add_shape(
                1, Inches(3.8), Inches(row_y + 0.02), Inches(0.55), Inches(0.22)
            )
            pct_box.fill.solid(); pct_box.fill.fore_color.rgb = GRAY
            pct_box.line.fill.background()
            pct_txb = slide.shapes.add_textbox(Inches(3.8), Inches(row_y + 0.02), Inches(0.55), Inches(0.22))
            pp2 = pct_txb.text_frame.paragraphs[0]
            pp2.alignment = 2  # center
            rr2 = pp2.add_run()
            rr2.text = _fmt_pct(t["pct"]); rr2.font.size = Pt(9)
            rr2.font.bold = True; rr2.font.color.rgb = WHITE; rr2.font.name = "Calibri"
            # Fecha
            tb_f = slide.shapes.add_textbox(Inches(4.42), Inches(row_y), Inches(1.4), Inches(0.26))
            pf   = tb_f.text_frame.paragraphs[0]
            rf   = pf.add_run()
            rf.text = t.get("fecha_vto", "") or ""
            rf.font.size = Pt(9.5); rf.font.color.rgb = GRAY; rf.font.name = "Calibri"

        # ── Nota VCP ─────────────────────────────────────────────────
        nota = slide.shapes.add_textbox(Inches(6.9), Inches(4.1), Inches(6.0), Inches(0.4))
        np_  = nota.text_frame.paragraphs[0]
        nr   = np_.add_run()
        nr.text = f"⚠ Variaciones VCP: completar en el HTML interactivo"
        nr.font.size = Pt(9); nr.font.color.rgb = RGBColor(0xb8, 0x5a, 0x1a)
        nr.font.name = "Calibri"

    prs.save(output_path)
    print(f"[generar_fondos] PPTX guardado: {output_path}")
    return output_path
